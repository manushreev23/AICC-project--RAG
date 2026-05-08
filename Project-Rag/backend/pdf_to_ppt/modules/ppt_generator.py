"""
ppt_generator.py
----------------
Generate a classroom-ready PPTX from summarized sections.

Design spec (per user requirements)
===================================
- Cream background (#FFF8E7)
- Soft orange / light-brown accents (#CC7722, #D2691E, #F4A261)
- Readable fonts (Calibri / Georgia)
- 5–10 slides minimum: Title, Intro, one per subsection (+ image), Examples, Conclusion
- Bullets, consistent spacing, alignment, no empty slides
"""

from __future__ import annotations

import os
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------- Theme constants ----------------------------- #

CREAM = RGBColor(0xFA, 0xF3, 0xE0)
ORANGE =RGBColor(0x2E, 0x4A, 0x62)   # primary accent
DARK_BROWN =  RGBColor(0xD9, 0x8C, 0x3F)     # body text
LIGHT_ORANGE =RGBColor(0x33, 0x33, 0x33)
TITLE_FONT = "Georgia"
BODY_FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# -------------------------- Low-level helpers ----------------------------- #

def _set_cream_background(slide) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = CREAM


def _add_accent_bar(slide) -> None:
    """Thin orange bar on the left edge for a consistent visual identity."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), SLIDE_H)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE


def _add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.name = BODY_FONT
    run.font.size = Pt(10)
    run.font.color.rgb = LIGHT_ORANGE
    run.font.italic = True


def _add_title(slide, text: str, top: float = 0.5) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(12.0), Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = TITLE_FONT
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = ORANGE

    # Thin underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(top + 1.05), Inches(1.5), Emu(30000)
    )
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_ORANGE


def _add_bullets(
    slide,
    bullets: List[str],
    left: float = 0.75,
    top: float = 1.9,
    width: float = 11.8,
    height: float = 4.8,
    size: int = 20,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = f"•  {b}"
        run.font.name = BODY_FONT
        run.font.size = Pt(size)
        run.font.color.rgb = DARK_BROWN


def _add_image(slide, image_path: str, left: float, top: float, width: float, height: float) -> bool:
    if not image_path or not os.path.exists(image_path):
        return False
    try:
        slide.shapes.add_picture(
            image_path, Inches(left), Inches(top), width=Inches(width), height=Inches(height)
        )
        return True
    except Exception:
        return False


# ------------------------------ Slide types -------------------------------- #

def _title_slide(prs: Presentation, topic: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_cream_background(slide)
    _add_accent_bar(slide)

    # Big centered title
    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.3), Inches(1.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = topic
    run.font.name = TITLE_FONT
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = ORANGE

    # Subtitle
    sub = slide.shapes.add_textbox(Inches(1.0), Inches(4.3), Inches(11.3), Inches(0.8))
    sp = sub.text_frame.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    srun = sp.add_run()
    srun.text = subtitle
    srun.font.name = BODY_FONT
    srun.font.size = Pt(22)
    srun.font.color.rgb = DARK_BROWN
    srun.font.italic = True

    # Decorative divider
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.8), Inches(5.2), Inches(1.7), Emu(40000)
    )
    div.line.fill.background()
    div.fill.solid()
    div.fill.fore_color.rgb = LIGHT_ORANGE


def _content_slide(
    prs: Presentation,
    title: str,
    bullets: List[str],
    image_path: Optional[str],
    footer: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_cream_background(slide)
    _add_accent_bar(slide)
    _add_title(slide, title)

    if image_path and os.path.exists(image_path):
        # Two-column layout: bullets left, image right.
        _add_bullets(slide, bullets, left=0.75, top=1.9, width=7.4, height=4.8, size=18)
        _add_image(slide, image_path, left=8.5, top=2.0, width=4.2, height=4.2)
    else:
        # Full-width bullets.
        _add_bullets(slide, bullets, left=0.75, top=1.9, width=11.8, height=4.8, size=20)

    _add_footer(slide, footer)


# ------------------------------ Public API --------------------------------- #

def generate_presentation(
    topic: str,
    intro_bullets: List[str],
    section_bullets: Dict[str, List[str]],
    conclusion_bullets: List[str],
    section_images: Optional[Dict[str, str]] = None,
    output_path: str = "output/presentation.pptx",
    min_slides: int = 5,
    max_slides: int = 10,
) -> str:
    """
    Build a PPTX and save to `output_path`. Returns the saved path.

    Slide plan:
      1. Title
      2. Introduction
      3..N. One slide per subsection (bullets + image)
      N+1. Examples / Key Takeaways (auto-generated from remaining content)
      Last. Conclusion
    """
    section_images = section_images or {}
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    footer = f"{topic} · Generated Presentation"

    # Count reserved slots dynamically (title is always added; intro/conclusion only if bullets exist)
    has_intro = bool(intro_bullets)
    has_conclusion = bool(conclusion_bullets)
    reserved = 1 + (1 if has_intro else 0) + (1 if has_conclusion else 0)

    # Budgets for body slides (subsections + optional Key Takeaways)
    body_budget_max = max(max_slides - reserved, 1)
    body_budget_min = max(min_slides - reserved, 0)

    # --- Slide 1: Title ---
    _title_slide(prs, topic, subtitle="An illustrated study overview")

    # --- Slide 2: Introduction ---
    if has_intro:
        _content_slide(prs, "Introduction", intro_bullets, section_images.get("__intro__"), footer)

    # --- Body slides: one per subsection, capped at body_budget_max ---
    ordered_sections = [(t, b) for t, b in section_bullets.items() if b]
    body_sections = ordered_sections[:body_budget_max]

    for section_title, bullets in body_sections:
        img = section_images.get(section_title)
        _content_slide(prs, section_title, bullets, img, footer)

    # --- Optional fill-in: Key Takeaways, only if body is below min budget AND room remains ---
    if len(body_sections) < body_budget_min and len(body_sections) < body_budget_max:
        takeaways = [b[0] for _, b in body_sections if b][:5]
        if takeaways:
            _content_slide(prs, "Key Takeaways", takeaways, None, footer)

    # --- Conclusion ---
    if has_conclusion:
        _content_slide(prs, "Conclusion", conclusion_bullets, section_images.get("__conclusion__"), footer)

    prs.save(output_path)
    return output_path


